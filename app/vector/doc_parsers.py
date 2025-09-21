from typing import Optional, List
import os

def read_text_any(path: str) -> Optional[str]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _read_pdf(path)
    if ext == ".docx":
        return _read_docx(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext in [".xlsx", ".xlsm"]:
        return _read_xlsx(path)
    if ext in [".md", ".txt"]:
        return open(path, "r", errors="ignore").read()
    if ext in [".html", ".htm"]:
        from bs4 import BeautifulSoup as BS
        html = open(path, "r", errors="ignore").read()
        return BS(html, "html.parser").get_text(" ")
    return None

def _read_pdf(path: str) -> Optional[str]:
    from pdfminer.high_level import extract_text
    try:
        return extract_text(path)
    except Exception:
        return None

def _read_docx(path: str) -> Optional[str]:
    try:
        from docx import Document
        d = Document(path)
        return "\n".join(p.text for p in d.paragraphs if p.text)
    except Exception:
        return None

def _read_pptx(path: str) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(p for p in parts if p)
    except Exception:
        return None

def _read_xlsx(path: str) -> Optional[str]:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts: List[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(v) for v in row if v is not None]
                if row_vals:
                    parts.append(" ".join(row_vals))
        return "\n".join(parts)
    except Exception:
        return None
